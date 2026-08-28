from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

NAVY = RGBColor(0x0E, 0x1B, 0x33)
AZUL = RGBColor(0x00, 0x55, 0xFF)
CORAL = RGBColor(0xFC, 0x60, 0x58)
CINZA = RGBColor(0x4A, 0x5A, 0x78)
CLARO = RGBColor(0xF4, 0xF7, 0xFC)
BRANCO = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, size=20, color=NAVY, bold=False, align=PP_ALIGN.LEFT,
          font="Calibri", anchor=MSO_ANCHOR.TOP, space=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    first = True
    for item in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6 * space)
        if isinstance(item, str):
            item = (item, {})
        txt, opts = item
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(opts.get("size", size))
        r.font.bold = opts.get("bold", bold)
        r.font.color.rgb = opts.get("color", color)
        r.font.name = opts.get("font", font)
    return tb


def header(slide, kicker, title):
    _rect(slide, 0, 0, prs.slide_width, Inches(1.25), NAVY)
    _rect(slide, Inches(0.45), Inches(0.33), Inches(0.09), Inches(0.6), CORAL)
    _text(slide, Inches(0.75), Inches(0.22), Inches(11), Inches(0.4),
          [(kicker, dict(size=13, color=RGBColor(0xA9,0xBC,0xDD), bold=True))])
    _text(slide, Inches(0.75), Inches(0.55), Inches(11.5), Inches(0.7),
          [(title, dict(size=30, color=BRANCO, bold=True))])


def footer(slide, n):
    _text(slide, Inches(0.45), Inches(7.05), Inches(5), Inches(0.35),
          [("Seazone · Hackathon Jovens Talentos AI Builder 2026", dict(size=10, color=CINZA))])
    _text(slide, Inches(12.3), Inches(7.05), Inches(0.7), Inches(0.35),
          [(str(n), dict(size=12, color=CINZA))], align=PP_ALIGN.RIGHT)


def bullets(slide, x, y, w, h, items, size=18, gap=10):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for it in items:
        if isinstance(it, tuple):
            head, body = it
        else:
            head, body = None, it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        p.level = 0
        if head:
            r = p.add_run()
            r.text = "• " + head + " "
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.color.rgb = NAVY
        if body:
            r = p.add_run()
            r.text = body
            r.font.size = Pt(size)
            r.font.color.rgb = CINZA


def tabela(slide, x, y, w, h, colunas, linhas, header_fill=NAVY, col_widths=None, fs=15):
    rows, cols = len(linhas) + 1, len(colunas)
    gt = slide.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = Emu(int(w * cw / total))
    for j, cname in enumerate(colunas):
        cell = gt.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        cell.text = cname
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.size = Pt(14)
                r.font.bold = True
                r.font.color.rgb = BRANCO
    for i, linha in enumerate(linhas, start=1):
        for j, val in enumerate(linha):
            cell = gt.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = CLARO if i % 2 else BRANCO
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(fs)
                    r.font.color.rgb = NAVY
    return gt


def build():
    # ===== Slide 1: Capa =====
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _rect(s, 0, Inches(6.1), prs.slide_width, Inches(1.4), AZUL)
    _text(s, Inches(0.9), Inches(1.2), Inches(11.5), Inches(0.5),
          [("HACKATHON · JOVENS TALENTOS AI BUILDER 2026", dict(size=14, color=RGBColor(0xA9,0xBC,0xDD), bold=True))])
    _text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.8),
          [("Recomendação de Investimento Imobiliário", dict(size=42, color=BRANCO, bold=True)),
           ("em Itapema (SC)", dict(size=42, color=CORAL, bold=True))])
    _text(s, Inches(0.9), Inches(3.7), Inches(11), Inches(0.6),
          [("Análise de mercado Airbnb + VivaReal | Decisão com apoio de IA", dict(size=18, color=RGBColor(0xD6,0xDF,0xF0)))] )
    _text(s, Inches(0.9), Inches(4.4), Inches(11.5), Inches(1.4),
          [("4 perguntas → 1 recomendação, com retorno estimado e posição sobre a tese dos compactos",
            dict(size=16, color=RGBColor(0xA9,0xBC,0xDD)))])
    _text(s, Inches(0.9), Inches(6.45), Inches(11.5), Inches(0.6),
          [("Relatório completo: relatorio.md  ·  Processo com IA: ai-log/  ·  Projeto: github.com/escalant-e/jt2026-ricardo-escalante",
            dict(size=14, color=BRANCO))])

    # ===== Slide 2: Contexto e pergunta =====
    s = prs.slides.add_slide(BLANK)
    header(s, "CONTEXTO", "O desafio")
    bullets(s, Inches(0.75), Inches(1.65), Inches(11.8), Inches(2.2), [
        ("Base de dados real de Itapema (SC):", "4.441 anúncios de Airbnb + 8.329 de venda (VivaReal)."),
        ("Missão:", "responder 4 perguntas e entregar uma recomendação de investimento para a Seazone, líder em short stay no Brasil."),
        ("Critérios 'abertos' de propósito:", "definimos o critério (retorno ajustado por ocupação) e justificamos com os dados."),
    ], size=19, gap=12)
    _rect(s, Inches(0.75), Inches(4.2), Inches(11.8), Inches(2.5), CLARO)
    _text(s, Inches(1.05), Inches(4.4), Inches(11.2), Inches(0.5),
          [("TESE INTERNA A AVALIAR", dict(size=13, color=AZUL, bold=True))])
    _text(s, Inches(1.05), Inches(4.8), Inches(11.2), Inches(1.7),
          [("\"Apartamentos compactos (studio/1Q) na região do Centro seriam a aposta mais eficiente para a Seazone.\"",
            dict(size=19, color=NAVY, bold=True)),
           ("→ Tomamos posição: os dados a sustentam ou não?", dict(size=17, color=CINZA))])
    footer(s, 2)

    # ===== Slide 3: Método e dados =====
    s = prs.slides.add_slide(BLANK)
    header(s, "MÉTODO", "Como analisamos")
    bullets(s, Inches(0.75), Inches(1.6), Inches(6.2), Inches(5), [
        ("Receita:", "diária mediana × 365 × ocupação (cenários 40/55/70%)."),
        ("Price_AV = disponibilidade", "→ modelamos por cenários, não 'dias ocupados'."),
        ("Retorno:", "receita ÷ preço de venda (VivaReal)."),
        ("Confiabilidade:", "pares ALTA/MÉDIA/BAIXA por volume de listings e anúncios."),
        ("Rigor:", "CRISP-DM + teste de falseamento da tese."),
    ], size=18, gap=12)
    _rect(s, Inches(7.3), Inches(1.6), Inches(5.3), Inches(5.1), NAVY)
    _text(s, Inches(7.6), Inches(1.85), Inches(4.7), Inches(0.5),
          [("BASES UTILIZADAS", dict(size=13, color=RGBColor(0xA9,0xBC,0xDD), bold=True))])
    dados = [("Details", "perfil, quartos, ratings"),
             ("Hosts", "anfitriões"),
             ("Mesh", "bairro / geolocalização"),
             ("Price_AV", "diária por data (jan–abr/25)"),
             ("VivaReal", "preço de venda, área, condomínio")]
    tb = s.shapes.add_textbox(Inches(7.6), Inches(2.35), Inches(4.7), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (a, b) in enumerate(dados):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        r = p.add_run(); r.text = f"{a} — "
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = BRANCO
        r2 = p.add_run(); r2.text = b
        r2.font.size = Pt(15); r2.font.color.rgb = RGBColor(0xD6,0xDF,0xF0)
    footer(s, 3)

    # ===== Slide 4: Q1 perfil em 3 níveis =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PERGUNTA 1", "Melhor perfil de imóvel (tipologia, nº de quartos)")
    bullets(s, Inches(0.75), Inches(1.55), Inches(11.8), Inches(0.9), [
        ("Definimos o critério:", "três óticas — receita bruta, produtividade de área (R$/m²) e retorno sobre capital."),
    ], size=17, gap=6)
    tabela(s, Inches(0.75), Inches(2.4), Inches(6.9), Inches(3.4),
           ["Tipologia", "Receita bruta (55%)", "R$/m²", "Retorno s/ capital"],
           [["4Q+", "R$ 210,8 mil", "R$ 1.114", "5,7%"],
            ["3Q", "R$ 130,5 mil", "R$ 1.183", "9,0%"],
            ["2Q", "R$ 90,3 mil", "R$ 1.202", "10,3%"],
            ["1Q", "R$ 77,3 mil", "R$ 1.807", "9,6%"]],
           col_widths=[1, 1.3, 1, 1.2], fs=15)
    _rect(s, Inches(7.95), Inches(2.4), Inches(4.6), Inches(1.4), CORAL)
    _text(s, Inches(8.2), Inches(2.6), Inches(4.1), Inches(1.1),
          [("3Q/4Q+ vencem", dict(size=16, color=BRANCO, bold=True)),
           ("no faturamento bruto", dict(size=13, color=BRANCO))])
    _rect(s, Inches(7.95), Inches(4.0), Inches(4.6), Inches(1.4), AZUL)
    _text(s, Inches(8.2), Inches(4.2), Inches(4.1), Inches(1.1),
          [("2Q é o equilíbrio", dict(size=16, color=BRANCO, bold=True)),
           ("melhor retorno sobre capital (10,3%)", dict(size=13, color=BRANCO))])
    _rect(s, Inches(7.95), Inches(5.6), Inches(4.6), Inches(1.4), NAVY)
    _text(s, Inches(8.2), Inches(5.8), Inches(4.1), Inches(1.1),
          [("1Q maximiza R$/m²", dict(size=16, color=BRANCO, bold=True)),
           ("(R$ 1.807) e tipo 'apartamento' domina", dict(size=13, color=BRANCO))])
    footer(s, 4)

    # ===== Slide 5: Q2 localização =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PERGUNTA 2", "Melhor localização em termos de receita")
    tabela(s, Inches(0.75), Inches(1.9), Inches(7.6), Inches(3.8),
           ["Bairro", "n", "Diária mediana", "Receita anual (55%)"],
           [["MEIA PRAIA", "630", "R$ 590", "R$ 118,4 mil"],
            ["CANTO DA PRAIA", "8", "R$ 559", "R$ 112,2 mil"],
            ["TABU. OLIVEIRAS", "20", "R$ 540", "R$ 108,4 mil"],
            ["CENTRO", "205", "R$ 509", "R$ 102,2 mil"],
            ["MORRETES", "82", "R$ 471", "R$ 94,5 mil"]],
           col_widths=[1.4, 0.5, 1, 1.1], fs=15)
    _rect(s, Inches(8.7), Inches(1.9), Inches(3.9), Inches(3.8), CLARO)
    _text(s, Inches(9.0), Inches(2.15), Inches(3.4), Inches(0.5),
          [("LEITURA", dict(size=13, color=AZUL, bold=True))])
    _text(s, Inches(9.0), Inches(2.55), Inches(3.4), Inches(3.0),
          [("Meia Praia lidera a receita consolidada.", dict(size=16, color=NAVY, bold=True)),
           ("Canto da Praia e Areal têm números altos, mas volume mínimo.", dict(size=14, color=CINZA)),
           ("Centro é a alternativa relevante de maior volume na região central.", dict(size=14, color=CINZA))])
    footer(s, 5)

    # ===== Slide 6: Q3 drivers =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PERGUNTA 3", "O que explica as melhores receitas?")
    bullets(s, Inches(0.75), Inches(1.6), Inches(6.5), Inches(3.6), [
        ("Regressão múltipla (R²=0,45):", "nº quartos, banheiros e hóspedes são os drivers mais fortes."),
        ("Ratings pouco relevantes:", "star_rating global não explica receita."),
        ("Cuidado estatístico:", "coeficientes negativos de ratings = colinearidade/mix de preço, não causalidade."),
        ("Conclusão:", "estrutura e capacidade do imóvel (tamanho) determinam a receita, não a avaliação."),
    ], size=18, gap=13)
    _rect(s, Inches(7.65), Inches(1.6), Inches(5.0), Inches(5.1), NAVY)
    _text(s, Inches(7.95), Inches(1.85), Inches(4.4), Inches(0.5),
          [("PRINCIPAIS DRIVERS", dict(size=13, color=RGBColor(0xA9,0xBC,0xDD), bold=True))])
    drivers = [("Quartos", "+21,1 ▲"), ("Banheiros", "+18,8 ▲"), ("Hóspedes", "+13,6 ▲"),
               ("Precisão do anúncio", "+22,8 ▲"), ("Satisfação", "+8,8 ▲")]
    tb = s.shapes.add_textbox(Inches(7.95), Inches(2.45), Inches(4.4), Inches(4.0))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (a, b) in enumerate(drivers):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(14)
        r = p.add_run(); r.text = a
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BRANCO
        r2 = p.add_run(); r2.text = "\t" + b
        r2.font.size = Pt(18); r2.font.color.rgb = CORAL
    footer(s, 6)

    # ===== Slide 7: Q4 decisão e retorno =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PERGUNTA 4", "O que comprar hoje — e por quê")
    tabela(s, Inches(0.75), Inches(1.75), Inches(7.6), Inches(2.9),
           ["Ativo", "Retorno Bruto 55%", "Retorno Liq. Encargos 55%"],
           [["MORRETES · 2Q", "11,3%", "10,8%"],
            ["CENTRO · 2Q", "9,8%", "9,2%"],
            ["MEIA PRAIA · 2Q", "8,5%", "7,9%"]],
           col_widths=[1.4, 1, 1.2], fs=16)
    bullets(s, Inches(0.75), Inches(5.0), Inches(7.6), Inches(2.2), [
        ("Por que 2Q em Morretes:", "melhor retorno (bruto 11,3%) com preço acessível (R$ 790 mil)."),
        ("Confiança:", "volumes grandes (59 listings / 1.043 anúncios) = ALTA confiabilidade."),
        ("Equilíbrio:", "liquidez, menor barreira de entrada e reprodutibilidade estatística."),
    ], size=16, gap=8)
    _rect(s, Inches(8.7), Inches(1.75), Inches(3.9), Inches(4.9), CORAL)
    _text(s, Inches(8.95), Inches(2.0), Inches(3.4), Inches(0.5),
          [("COMO COMPARAR", dict(size=13, color=BRANCO, bold=True))])
    _text(s, Inches(8.95), Inches(2.4), Inches(3.4), Inches(4.1),
          [("Receita bruta = diária mediana × 365 × ocupação.", dict(size=15, color=BRANCO)),
           ("Retr. líquido = receita − condomínio×12 − IPTU.", dict(size=15, color=BRANCO)),
           ("3 cenários: 40% / 55% / 70%.", dict(size=15, color=BRANCO)),
           ("Não conta opex (limpeza, energia, taxa de plataforma).", dict(size=13, color=BRANCO))])
    footer(s, 7)

    # ===== Slide 8: Tese dos compactos =====
    s = prs.slides.add_slide(BLANK)
    header(s, "POSIÇÃO SOBRE A TESE", "Compactos (1Q) no Centro — os dados sustentam?")
    _rect(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(1.2), CLARO)
    _text(s, Inches(1.1), Inches(1.9), Inches(11.1), Inches(0.9),
          [("Resposta: ", dict(size=22, color=NAVY, bold=True)),
           ("parcialmente validada e reposicionada — não refutada.", dict(size=22, color=AZUL, bold=True))])
    bullets(s, Inches(0.75), Inches(3.2), Inches(11.8), Inches(3.6), [
        ("O que confirma:", "unidades pequenas têm alta eficiência de capital (receita/m² e retorno)."),
        ("O que ajusta:", "o ótimo não é o 1Q, e sim o 2Q (10,3% vs 9,6% do 1Q)."),
        ("Onde:", "o melhor retorno confiável está em Morretes 2Q (11,3%), não no Centro 1Q."),
        ("Decisão:", "investir em 2Q (não 1Q), priorizar Morretes e depois Centro; Meia Praia para receita."),
    ], size=19, gap=14)
    footer(s, 8)

    # ===== Slide 9: Uso da IA =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PROCESSO", "Como a IA foi usada")
    bullets(s, Inches(0.75), Inches(1.6), Inches(11.8), Inches(4.4), [
        ("Ferramenta colaborativa, não substituta:", "a IA testou hipóteses e o analista decidiu cada critério."),
        ("Ponto de virada (Price_AV):", "diagnóstico mostrou disponibilidade, não ocupação → modelamos por cenários."),
        ("Qualidade:", "detectou inconsistências (owner duplicado, n pequeno, colinearidade) e auditou o relatório como revisor sênior."),
        ("Decisões do analista:", "limiar de R$ 3.000, taxas 40/55/70%, filtros do VivaReal, nomenclatura."),
        ("Rastreabilidade:", "tudo em ai-log/ (conversa completa em texto) e scripts 01–15."),
    ], size=18, gap=12)
    footer(s, 9)

    # ===== Slide 10: Próximos passos / mais uma semana =====
    s = prs.slides.add_slide(BLANK)
    header(s, "PRÓXIMOS PASSOS", "Com mais uma semana")
    bullets(s, Inches(0.75), Inches(1.7), Inches(11.8), Inches(4.6), [
        ("Ano completo:", "calibrar sazonalidade real e ocupação/inverno (remover a maior limitação)."),
        ("Custos operacionais:", "transformar 'líquido de encargos' em NOI/cap rate completo."),
        ("Modelo mais robusto:", "validação cruzada, VIF e variáveis derivadas (distância à praia, oferta)."),
        ("Benchmark Seazone + validação externa:", "testar a hipótese de renda da terra em Morretes com dados censitários."),
        ("Sensibilidade:", "testar o limiar de confiabilidade (5/10/15/20) para blindar a decisão 2Q."),
    ], size=19, gap=14)
    footer(s, 10)

    # ===== Slide 11: Conclusão =====
    s = prs.slides.add_slide(BLANK)
    _rect(s, 0, 0, prs.slide_width, prs.slide_height, NAVY)
    _rect(s, 0, Inches(6.2), prs.slide_width, Inches(1.3), AZUL)
    _text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.5),
          [("CONCLUSÃO", dict(size=14, color=RGBColor(0xA9,0xBC,0xDD), bold=True))])
    _text(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(2.2),
          [("Investir em apartamentos 2Q,", dict(size=40, color=BRANCO, bold=True)),
           ("começando por Morretes.", dict(size=40, color=CORAL, bold=True))])
    _text(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(2.2),
          [("Retorno líquido de encargos de 10,8% a. a. (cenário base de ocupação 55%), "
            "com ticket acessível e alta confiabilidade.",
            dict(size=20, color=RGBColor(0xD6,0xDF,0xF0))),
           ("Tese dos compactos: validada na intuição (eficiência), ajustada na execução (2Q, Morretes).",
            dict(size=18, color=RGBColor(0xA9,0xBC,0xDD)))])
    _text(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.7),
          [("Obrigado!", dict(size=22, color=BRANCO, bold=True)),
           ("   Relatório: relatorio.md · Código: scripts/ · Processo: ai-log/",
            dict(size=14, color=BRANCO))])
    footer(s, 11)

    out = os.path.join(os.path.dirname(__file__), "..", "output", "apresentacao_seazone_itapema.pptx")
    prs.save(out)
    print("Salvo:", out, "(", len(prs.slides.__iter__.__self__._sldIdLst), "slides )")


import os
if __name__ == "__main__":
    build()